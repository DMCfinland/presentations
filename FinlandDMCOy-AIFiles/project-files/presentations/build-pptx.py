#!/usr/bin/env python3
"""
Finland DMC — Presentation Generator
Builds 5 PPTX presentations following Levi Tour template structure.
Brand: FIN+LAND DMC (blue #0035A0, gold #C9A227, cream backgrounds)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ========== BRAND CONSTANTS ==========
BRAND_BLUE = RGBColor(0x00, 0x35, 0xA0)
BRAND_BLUE_LIGHT = RGBColor(0x1A, 0x5B, 0xC4)
STONE_DARK = RGBColor(0x1A, 0x26, 0x34)
STONE_MEDIUM = RGBColor(0x2D, 0x3E, 0x50)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
CREAM = RGBColor(0xF5, 0xF0, 0xEB)
CREAM_DARK = RGBColor(0xED, 0xE6, 0xDD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x2C, 0x2C, 0x2C)
TEXT_MEDIUM = RGBColor(0x5A, 0x5A, 0x5A)
TEXT_LIGHT = RGBColor(0x8A, 0x8A, 0x8A)

SLIDE_W = Inches(13.333)  # 16:9 widescreen
SLIDE_H = Inches(7.5)

LOGO_PATH = os.path.expanduser("~/Downloads/Media/Fin+Land-dmc.png")
OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ========== HELPER FUNCTIONS ==========

def new_prs():
    """Create a new 16:9 presentation with cream background default."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def add_blank_slide(prs):
    """Add a blank slide."""
    layout = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(layout)

def set_bg_color(slide, color):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, color, left=0, top=0, width=None, height=None):
    """Add a colored rectangle as background."""
    w = width or SLIDE_W
    h = height or SLIDE_H
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, text, left, top, width, height,
                 font_name='Montserrat', font_size=18, bold=False, italic=False,
                 color=TEXT_DARK, alignment=PP_ALIGN.LEFT, font_family_fallback='Arial'):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox

def add_multi_text(slide, texts, left, top, width, height):
    """Add text box with multiple paragraphs. texts = [(text, size, bold, italic, color, align), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, italic, color, align) in enumerate(texts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.name = 'Montserrat'
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        p.alignment = align
        p.space_after = Pt(8)
    return txBox

def add_logo(slide, right_margin=Inches(0.5), bottom_margin=Inches(0.3), height=Inches(0.5)):
    """Add FIN+LAND logo to bottom-right corner."""
    if os.path.exists(LOGO_PATH):
        left = SLIDE_W - Inches(2.5) - right_margin
        top = SLIDE_H - height - bottom_margin
        slide.shapes.add_picture(LOGO_PATH, left, top, height=height)

def add_photo_placeholder(slide, left, top, width, height, label="[ADD PHOTO]"):
    """Add a labeled placeholder box where a photo should go."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFC)  # ice-light
    shape.line.color.rgb = RGBColor(0xB8, 0xD4, 0xE8)  # ice-blue
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = 'Montserrat'
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape

def add_stat(slide, number, label, left, top, num_size=48, label_size=12, num_color=BRAND_BLUE):
    """Add a stat block (big number + small label)."""
    add_text_box(slide, number, left, top, Inches(3), Inches(0.8),
                 font_name='Playfair Display', font_size=num_size, bold=True,
                 color=num_color)
    add_text_box(slide, label, left, top + Inches(0.75), Inches(3), Inches(0.4),
                 font_size=label_size, color=TEXT_MEDIUM)

# ========== SLIDE BUILDERS ==========

def build_cover(prs, title, subtitle, photo_label="[COVER PHOTO: Full-bleed landscape]"):
    """T1: Full-color cover slide."""
    slide = add_blank_slide(prs)
    add_shape_bg(slide, STONE_DARK)
    add_photo_placeholder(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, photo_label)
    # Overlay
    overlay = add_shape_bg(slide, STONE_DARK, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = STONE_DARK
    # Title
    add_text_box(slide, title, Inches(1), Inches(2), Inches(11), Inches(2),
                 font_name='Playfair Display', font_size=60, bold=True,
                 color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, subtitle, Inches(2), Inches(4.2), Inches(9), Inches(1),
                 font_name='Playfair Display', font_size=24, italic=True,
                 color=RGBColor(0xE8, 0xF4, 0xFC), alignment=PP_ALIGN.CENTER)
    # Brand
    add_text_box(slide, "FINLAND DMC", Inches(4.5), Inches(6.5), Inches(4), Inches(0.5),
                 font_size=12, bold=True, color=RGBColor(0x7E, 0xB8, 0xD8),
                 alignment=PP_ALIGN.CENTER)
    add_logo(slide)
    return slide

def build_stats(prs, title, stats, photo_label="[PHOTO: Destination hero shot]"):
    """T2: Welcome/Stats slide. stats = [(number, label), ...]"""
    slide = add_blank_slide(prs)
    set_bg_color(slide, CREAM)
    # Photo left side
    add_photo_placeholder(slide, Inches(0), Inches(0), Inches(5), SLIDE_H, photo_label)
    # Title
    add_text_box(slide, title, Inches(5.5), Inches(0.8), Inches(7), Inches(0.8),
                 font_name='Playfair Display', font_size=36, bold=True, color=STONE_DARK)
    # Stats grid
    for i, (num, lbl) in enumerate(stats):
        col = i % 2
        row = i // 2
        x = Inches(5.5) + col * Inches(3.5)
        y = Inches(2.0) + row * Inches(1.5)
        add_stat(slide, num, lbl, x, y)
    add_logo(slide)
    return slide

def build_content_lr(prs, title, paragraphs, photo_label="[PHOTO]"):
    """T3: Text left, photo right."""
    slide = add_blank_slide(prs)
    set_bg_color(slide, CREAM)
    # Title
    add_text_box(slide, title, Inches(0.8), Inches(0.8), Inches(6), Inches(0.8),
                 font_name='Playfair Display', font_size=32, bold=True, color=STONE_DARK)
    # Paragraphs
    texts = [(p, 14, False, False, TEXT_MEDIUM, PP_ALIGN.LEFT) for p in paragraphs]
    add_multi_text(slide, texts, Inches(0.8), Inches(2.0), Inches(6), Inches(4.5))
    # Photo right
    add_photo_placeholder(slide, Inches(7.5), Inches(0.3), Inches(5.5), Inches(6.8), photo_label)
    add_logo(slide)
    return slide

def build_content_rl(prs, title, paragraphs, photo_labels=None):
    """T4: Photos left, text right."""
    slide = add_blank_slide(prs)
    set_bg_color(slide, CREAM)
    # Photos left (2 stacked)
    labels = photo_labels or ["[PHOTO 1]", "[PHOTO 2]"]
    add_photo_placeholder(slide, Inches(0.3), Inches(0.3), Inches(5), Inches(3.4), labels[0])
    add_photo_placeholder(slide, Inches(0.3), Inches(3.8), Inches(5), Inches(3.4), labels[1])
    # Title + text right
    add_text_box(slide, title, Inches(5.8), Inches(0.8), Inches(7), Inches(1),
                 font_name='Playfair Display', font_size=30, bold=True, color=STONE_DARK)
    texts = [(p, 14, False, False, TEXT_MEDIUM, PP_ALIGN.LEFT) for p in paragraphs]
    add_multi_text(slide, texts, Inches(5.8), Inches(2.2), Inches(6.8), Inches(4.5))
    add_logo(slide)
    return slide

def build_activities(prs, title, subtitle, activities, photo_labels=None):
    """T5: Activities with photo grid."""
    slide = add_blank_slide(prs)
    set_bg_color(slide, CREAM)
    # Title + subtitle
    add_text_box(slide, title, Inches(0.8), Inches(0.5), Inches(6), Inches(1),
                 font_name='Playfair Display', font_size=28, bold=True, color=STONE_DARK)
    add_text_box(slide, subtitle, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
                 font_size=12, italic=True, color=TEXT_MEDIUM)
    # Activity list in 2 columns
    col1 = activities[:len(activities)//2]
    col2 = activities[len(activities)//2:]
    texts1 = [("  \u2022  " + a, 12, False, False, TEXT_DARK, PP_ALIGN.LEFT) for a in col1]
    texts2 = [("  \u2022  " + a, 12, False, False, TEXT_DARK, PP_ALIGN.LEFT) for a in col2]
    add_multi_text(slide, texts1, Inches(0.8), Inches(2.2), Inches(3), Inches(4.5))
    add_multi_text(slide, texts2, Inches(3.8), Inches(2.2), Inches(3), Inches(4.5))
    # Photo grid (2x3) on right
    labels = photo_labels or [f"[ACTIVITY PHOTO {i+1}]" for i in range(6)]
    for i in range(6):
        col = i % 2
        row = i // 2
        x = Inches(7.3) + col * Inches(3)
        y = Inches(0.3) + row * Inches(2.4)
        add_photo_placeholder(slide, x, y, Inches(2.9), Inches(2.3), labels[min(i, len(labels)-1)])
    add_logo(slide)
    return slide

def build_destinations(prs, title, subtitle, destinations):
    """Destination cards overview. destinations = [(name, tagline, details), ...]"""
    slide = add_blank_slide(prs)
    set_bg_color(slide, CREAM)
    add_text_box(slide, title, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 font_name='Playfair Display', font_size=32, bold=True, color=STONE_DARK)
    add_text_box(slide, subtitle, Inches(0.8), Inches(1.3), Inches(10), Inches(0.5),
                 font_size=13, italic=True, color=TEXT_MEDIUM)
    # Destination cards
    card_w = Inches(11.5) / len(destinations) - Inches(0.2)
    for i, (name, tagline, details) in enumerate(destinations):
        x = Inches(0.8) + i * (card_w + Inches(0.3))
        y = Inches(2.2)
        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xFA, 0xF8, 0xF5)
        card.line.color.rgb = CREAM_DARK
        card.line.width = Pt(1)
        # Photo placeholder in card
        add_photo_placeholder(slide, x + Inches(0.15), y + Inches(0.15),
                            card_w - Inches(0.3), Inches(2), f"[PHOTO: {name}]")
        # Name
        add_text_box(slide, name, x + Inches(0.2), y + Inches(2.3), card_w - Inches(0.4), Inches(0.5),
                     font_name='Playfair Display', font_size=20, bold=True, color=BRAND_BLUE,
                     alignment=PP_ALIGN.CENTER)
        # Tagline
        add_text_box(slide, tagline, x + Inches(0.2), y + Inches(2.8), card_w - Inches(0.4), Inches(0.4),
                     font_size=11, italic=True, color=TEXT_MEDIUM, alignment=PP_ALIGN.CENTER)
        # Details
        add_text_box(slide, details, x + Inches(0.2), y + Inches(3.3), card_w - Inches(0.4), Inches(1.3),
                     font_size=10, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
    add_logo(slide)
    return slide

def build_product(prs, title, subtitle, nights):
    """T8: Product examples with night options. nights = [(n, includes), ...]"""
    slide = add_blank_slide(prs)
    set_bg_color(slide, CREAM)
    add_text_box(slide, title, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 font_name='Playfair Display', font_size=30, bold=True, color=STONE_DARK)
    add_text_box(slide, subtitle, Inches(0.8), Inches(1.3), Inches(10), Inches(0.6),
                 font_size=13, italic=True, color=TEXT_MEDIUM)
    # Night cards
    card_w = Inches(11.5) / len(nights) - Inches(0.15)
    for i, (n, includes) in enumerate(nights):
        x = Inches(0.8) + i * (card_w + Inches(0.2))
        y = Inches(2.3)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xFA, 0xF8, 0xF5)
        card.line.color.rgb = CREAM_DARK
        card.line.width = Pt(1)
        # Night number
        add_text_box(slide, str(n), x + Inches(0.1), y + Inches(0.3), card_w - Inches(0.2), Inches(0.7),
                     font_name='Playfair Display', font_size=40, bold=True, color=BRAND_BLUE,
                     alignment=PP_ALIGN.CENTER)
        label = "night" if n == 1 else "nights"
        add_text_box(slide, label, x + Inches(0.1), y + Inches(1.0), card_w - Inches(0.2), Inches(0.3),
                     font_size=11, color=TEXT_MEDIUM, alignment=PP_ALIGN.CENTER)
        # Includes
        add_text_box(slide, includes, x + Inches(0.15), y + Inches(1.6), card_w - Inches(0.3), Inches(2.5),
                     font_size=10, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
    add_logo(slide)
    return slide

def build_closing(prs, text, photo_label="[CLOSING: Full-bleed atmospheric photo]"):
    """T6: Impact/closing slide."""
    slide = add_blank_slide(prs)
    add_shape_bg(slide, STONE_DARK)
    add_photo_placeholder(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, photo_label)
    # Overlay
    add_shape_bg(slide, STONE_DARK)
    # Closing text
    add_text_box(slide, text, Inches(1.5), Inches(2.5), Inches(10), Inches(2.5),
                 font_name='Playfair Display', font_size=36, italic=True,
                 color=WHITE, alignment=PP_ALIGN.CENTER)
    add_logo(slide, bottom_margin=Inches(0.5))
    return slide

def build_company(prs):
    """T7: Company/contact slide."""
    slide = add_blank_slide(prs)
    set_bg_color(slide, CREAM)
    # Logo
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(4.5), Inches(0.5), height=Inches(1))
    # Brand name fallback
    add_text_box(slide, "FIN+LAND", Inches(3), Inches(0.6), Inches(7), Inches(0.8),
                 font_name='Playfair Display', font_size=36, bold=True, color=BRAND_BLUE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, "DESTINATION MANAGEMENT COMPANY", Inches(3), Inches(1.4), Inches(7), Inches(0.4),
                 font_size=11, color=TEXT_MEDIUM, alignment=PP_ALIGN.CENTER)
    # Tagline
    add_text_box(slide, "We know hospitality.", Inches(3), Inches(2.2), Inches(7), Inches(0.5),
                 font_name='Playfair Display', font_size=20, italic=True, color=TEXT_MEDIUM,
                 alignment=PP_ALIGN.CENTER)
    # Description
    desc = (
        "Finland DMC is a specialist destination management company covering all of Finland "
        "— from the Arctic wilderness of Lapland to the thousand lakes of Saimaa "
        "and the Nordic design capital of Helsinki.\n\n"
        "For over a decade, we have been crafting seamless, end-to-end travel experiences "
        "for international tour operators. From carefully selected accommodation and dining "
        "to inspiring activities and excursions, we bring every element together into one "
        "smooth and memorable whole.\n\n"
        "We serve tour operators, corporate groups, and event organisers worldwide. "
        "With deep local knowledge, long-term supplier partnerships, and a genuine passion "
        "for Finland, we turn trips into meaningful experiences."
    )
    add_text_box(slide, desc, Inches(2.5), Inches(3.0), Inches(8), Inches(3),
                 font_size=12, color=TEXT_MEDIUM, alignment=PP_ALIGN.CENTER)
    # Contact
    contact = "Finland DMC Oy  |  www.finland-dmc.com  |  info@finland-dmc.com  |  +358 43 2010 687"
    add_text_box(slide, contact, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
                 font_size=11, bold=True, color=STONE_DARK, alignment=PP_ALIGN.CENTER)
    return slide


# ========================================================================
# PRESENTATION 1: FINLAND — Country Overview (6 pages)
# ========================================================================
def build_01_finland():
    prs = new_prs()

    # 1. Cover
    build_cover(prs, "Finland", "The Land of a Thousand Experiences",
                "[COVER: Aerial view of Finnish landscape — lakes, forests, Northern Lights]")

    # 2. Welcome to Finland
    build_stats(prs, "Welcome to Finland!", [
        ("3", "unique regions: Lapland, Lakeland, Helsinki"),
        ("200K+", "lakes — the Land of a Thousand Lakes"),
        ("4", "distinct seasons"),
        ("#1", "happiest country in the world (2018–2026)"),
        ("Sep–Apr", "winter season with Northern Lights"),
    ], "[PHOTO: Finland landscape — summer lakeland or winter Lapland]")

    # 3. Lapland
    build_content_lr(prs, "Lapland — The Arctic North",
        ["Finnish Lapland is a land of extremes — from the midnight sun of summer "
         "to the magical polar nights of winter, illuminated by the Northern Lights.",
         "Home to world-class ski resorts, wilderness safaris, and Santa Claus himself, "
         "Lapland offers unforgettable Arctic experiences year-round.",
         "Our destinations: Levi, Saariselkä, and Pyhä — each with its own character, "
         "all with the magic of the Arctic."],
        "[PHOTO: Lapland winter — snow-covered trees, aurora, or fell landscape]")

    # 4. Lakeland
    build_content_lr(prs, "Lakeland — The Heart of Finland",
        ["The Finnish Lakeland is Europe's largest lake district — a breathtaking "
         "landscape of thousands of lakes, islands, and endless forests.",
         "In summer, it offers midnight sun lake cruises, fishing, hiking, and the "
         "famous Finnish sauna culture at its most authentic.",
         "Our destinations: Järvisydän (Rantasalmi), Tahko, and Sahalahti — "
         "where nature, luxury, and Finnish tradition meet."],
        "[PHOTO: Lakeland summer — aerial lake view, sunset over water, or Järvisydän resort]")

    # 5. Helsinki
    build_content_lr(prs, "Helsinki — The Nordic Capital",
        ["Helsinki is a vibrant, compact city where Nordic design meets world-class dining, "
         "fascinating architecture, and a unique island-dotted coastline.",
         "From the iconic Helsinki Cathedral to the cutting-edge Oodi Library, "
         "from market hall food tours to archipelago boat trips — Helsinki packs "
         "an extraordinary amount into a walkable city.",
         "The perfect gateway to start or end any Finland experience."],
        "[PHOTO: Helsinki — Cathedral, Senate Square, or waterfront panorama]")

    # 6. Map + Company
    slide = add_blank_slide(prs)
    set_bg_color(slide, CREAM)
    add_photo_placeholder(slide, Inches(0.5), Inches(0.5), Inches(5.5), Inches(6.5),
                         "[MAP OF FINLAND: Show Lapland (north), Lakeland (center-east), Helsinki (south)\n"
                         "Mark destinations: Levi, Saariselkä, Pyhä, Rantasalmi/Järvisydän, Tahko, Sahalahti, Helsinki]")
    add_text_box(slide, "Three Regions,\nOne Finland", Inches(6.5), Inches(1), Inches(6), Inches(1.5),
                 font_name='Playfair Display', font_size=32, bold=True, color=STONE_DARK)
    add_text_box(slide,
        "From the Arctic wilderness of Lapland to the thousand lakes "
        "of Saimaa and the Nordic design capital of Helsinki — "
        "Finland DMC covers it all.\n\n"
        "We create seamless multi-destination itineraries "
        "that showcase the best of Finland.",
        Inches(6.5), Inches(3), Inches(6), Inches(3),
        font_size=14, color=TEXT_MEDIUM)
    add_text_box(slide, "www.finland-dmc.com", Inches(6.5), Inches(6.2), Inches(6), Inches(0.5),
                 font_size=14, bold=True, color=BRAND_BLUE)
    add_logo(slide)

    path = os.path.join(OUTPUT_DIR, "01-Finland-Overview.pptx")
    prs.save(path)
    print(f"  Saved: {path}")
    return path


# ========================================================================
# PRESENTATION 2: FINLAND DMC — Company Intro (1 page)
# ========================================================================
def build_02_dmc():
    prs = new_prs()
    build_company(prs)
    path = os.path.join(OUTPUT_DIR, "02-Finland-DMC.pptx")
    prs.save(path)
    print(f"  Saved: {path}")
    return path


# ========================================================================
# PRESENTATION 3: LAPLAND — Winter (11 pages)
# ========================================================================
def build_03_lapland():
    prs = new_prs()

    # 1. Cover
    build_cover(prs, "Lapland", "The Arctic Winter Experience\nin Finnish Lapland",
                "[COVER: Northern Lights over snow-covered trees in Lapland]")

    # 2. Stats
    build_stats(prs, "Welcome to Lapland!", [
        ("3", "unique winter destinations"),
        ("200+", "Northern Lights nights per year"),
        ("50+", "winter activities available"),
        ("-30\u00b0C to +5\u00b0C", "winter temperature range"),
        ("Sep\u2013Apr", "winter season"),
    ], "[PHOTO: Lapland fell landscape at golden hour]")

    # 3. Northern Lights
    build_content_lr(prs, "The Northern Lights",
        ["Finnish Lapland is one of the best places in the world to experience "
         "the Aurora Borealis. Located far above the Arctic Circle and directly "
         "under the auroral oval, Lapland offers exceptional chances to see the "
         "Northern Lights on clear nights.",
         "Thanks to long polar nights, pristine Arctic air, and very low light "
         "pollution, the skies across our three destinations are ideal for "
         "aurora watching from September through March.",
         "We offer dedicated Northern Lights hunts by snowmobile, reindeer "
         "sleigh, and on foot — led by expert guides who know the best viewing "
         "spots and conditions."],
        "[PHOTO: Aurora Borealis over snowy landscape]")

    # 4. Destinations overview
    build_destinations(prs, "Three Winter Destinations",
        "Each with its own character, all with the magic of the Arctic",
        [("Levi", "Lapland's Largest Resort",
          "#1 ski resort in Lapland\n111 Northern Lights sightings/year\n60 restaurants & bars\n25,000 beds\n750,000 visitors/year"),
         ("Saariselkä", "Gateway to the Wilderness",
          "Urho Kekkonen National Park\nGlass igloos & aurora cabins\nGold panning heritage\nFell hiking & skiing"),
         ("Pyhä", "Untouched & Authentic",
          "Pyhä-Luosto National Park\nFinland's oldest ski resort\nAmethyst mine visits\nSmall-scale, genuine Arctic")])

    # 5. Activities
    build_activities(prs, "Things to do in the\nArctic Winter!",
        "To offer the best possible experience, we have a wealth of activities:",
        ["Northern Lights hunting", "Husky safari", "Reindeer sleigh ride",
         "Snowmobile safari", "Ice fishing", "Snowshoeing",
         "Cross-country skiing", "Alpine skiing", "Ice karting",
         "Traditional Finnish sauna", "Ice swimming", "Snow Castle visit",
         "Helicopter & snowcat tours", "Meet Santa Claus"],
        ["[Husky safari]", "[Snowmobile]", "[Reindeer]",
         "[Ice fishing]", "[Sauna & ice swim]", "[Northern Lights]"])

    # 6. Hotels
    build_content_rl(prs, "Accommodation — Hotels",
        ["Our Lapland destinations offer a perfect balance of Arctic atmosphere "
         "and modern comfort. From cosy alpine hotels to stylish design stays "
         "and family-friendly spa hotels, there's something for every traveller.",
         "Hotels are located close to the slopes, activities, and village services, "
         "ensuring easy access to everything Lapland has to offer.",
         "We work with all the leading hotel partners across Levi, Saariselkä, "
         "and Pyhä to find the perfect match for your clients."],
        ["[PHOTO: Lapland hotel exterior — winter]", "[PHOTO: Lapland hotel interior or lobby]"])

    # 7. Cabins
    build_content_rl(prs, "Accommodation — Private Cabins & Chalets",
        ["For those seeking privacy and authenticity, Lapland offers exceptional "
         "private cabins and chalets — from luxury villas with private saunas to "
         "cosy wilderness cabins surrounded by forest.",
         "Whether your clients want complete privacy or easy access to the village "
         "centre, this accommodation type provides flexibility tailored to any "
         "group size and preference."],
        ["[PHOTO: Luxury cabin exterior — snow]", "[PHOTO: Cabin interior — fireplace, sauna]"])

    # 8. Igloos
    build_content_rl(prs, "Accommodation — Northern Lights Igloos",
        ["The Aurora Cabins and Northern Lights Igloos are a magical way to "
         "experience the wintry star-filled sky and the Northern Lights while "
         "lying comfortably in a warm bed.",
         "Lapland has multiple resorts with glass-roofed accommodation, located "
         "away from light pollution to offer a peaceful stay and optimal setting "
         "for Northern Lights viewing.",
         "These unique stays are among the most sought-after experiences in all of Finland."],
        ["[PHOTO: Glass igloo exterior — aerial view]", "[PHOTO: Inside glass igloo — aurora visible]"])

    # 9. Product examples
    build_product(prs, "Build Your Lapland Experience",
        "Top things to do — your clients choose how many days and what they want",
        [(2, "Quick Arctic escape.\nNorthern Lights hunt\n+ husky safari\n+ hotel"),
         (3, "Classic Lapland.\nReindeer + snowmobile\n+ Northern Lights\n+ glass igloo option"),
         (4, "Deep Arctic.\nAdd ice fishing,\nsnowshoeing, sauna\n& ice swim ritual"),
         (5, "Multi-destination.\nCombine Levi +\nSaariselkä or add\nPyhä amethyst day"),
         (6, "Full experience.\nTwo destinations +\nwilderness cabin night\n+ all highlights"),
         (7, "Complete Arctic journey.\nThree destinations,\nevery experience,\nunforgettable")])

    # 10. Closing
    build_closing(prs, "All of this and much more\n— Lapland awaits!",
                  "[CLOSING: Full-bleed Northern Lights or snowy landscape]")

    # 11. Company
    build_company(prs)

    path = os.path.join(OUTPUT_DIR, "03-Lapland-Winter.pptx")
    prs.save(path)
    print(f"  Saved: {path}")
    return path


# ========================================================================
# PRESENTATION 4: LAKELAND — Summer (11 pages)
# ========================================================================
def build_04_lakeland():
    prs = new_prs()

    # 1. Cover
    build_cover(prs, "Lakeland", "The Arctic Summer Experience\nin the Finnish Lake District",
                "[COVER: Aerial view of Finnish lakeland — islands, forests, midnight sun]")

    # 2. Stats
    build_stats(prs, "Welcome to Lakeland!", [
        ("3", "unique summer destinations"),
        ("188,000", "lakes in Finland"),
        ("73", "days of midnight sun"),
        ("+30\u00b0C", "peak summer temperature"),
        ("May\u2013Sep", "summer season"),
    ], "[PHOTO: Lakeland aerial — lake, islands, sunset]")

    # 3. The Midnight Sun
    build_content_lr(prs, "The Midnight Sun",
        ["The Finnish Lakeland is one of the most magical places on Earth during "
         "summer. The sun barely sets, creating endless golden evenings perfect "
         "for lakeside dining, midnight fishing, and aurora-like pink skies.",
         "Europe's largest lake district stretches across central and eastern "
         "Finland — a breathtaking mosaic of lakes, islands, and ancient forests "
         "that offers peace, adventure, and authentic Finnish culture.",
         "From lakeside saunas to wilderness hikes, from berry picking to "
         "boat cruises — the Finnish summer is an experience unlike any other."],
        "[PHOTO: Midnight sun over lake — golden light, reflections]")

    # 4. Destinations
    build_destinations(prs, "Three Summer Destinations",
        "Each with its own character, all with the magic of a thousand lakes",
        [("Järvisydän\nRantasalmi", "Heart of the Lake",
          "Luxury lakeside resort\nSaimaa ringed seal habitat\n20 years of ice trails (winter)\nNational park gateway\nSpa & wellness"),
         ("Tahko", "Active Adventure Hub",
          "Year-round resort\nGolf, hiking, MTB trails\nLake Syväri activities\nFamily-friendly\nClose to Kuopio city"),
         ("Sahalahti", "Hidden Gem of Lakeland",
          "Authentic Finnish countryside\nPrivate lake houses\nBerry & mushroom foraging\nPeaceful retreats\nCultural experiences")])

    # 5. Activities
    build_activities(prs, "Things to do in the\nFinnish Summer!",
        "To offer the best possible experience, we have a wealth of activities:",
        ["Lake cruises & boat trips", "Finnish sauna experience", "Midnight sun fishing",
         "Hiking & nature walks", "Berry & mushroom picking", "Canoeing & kayaking",
         "Stand-up paddleboarding", "Cycling & mountain biking", "Swimming in pristine lakes",
         "Saimaa ringed seal watching", "Local farm visits", "Smoke sauna ritual",
         "Archipelago island hopping", "Finnish cooking workshops"],
        ["[Lake cruise]", "[Sauna by lake]", "[Midnight fishing]",
         "[Forest hiking]", "[Kayaking]", "[Sunset over lake]"])

    # 6. Hotels
    build_content_rl(prs, "Accommodation — Lakeside Hotels & Resorts",
        ["The Finnish Lakeland offers unique accommodation with stunning lake views. "
         "From the award-winning Järvisydän resort with its spa and fine dining "
         "to charming boutique hotels surrounded by nature.",
         "Most hotels offer direct lake access, saunas, and a wide range of "
         "activities right from the doorstep.",
         "We partner with the finest accommodation providers to ensure your "
         "clients experience the best of Finnish lakeside living."],
        ["[PHOTO: Järvisydän resort or lakeside hotel]", "[PHOTO: Hotel restaurant or spa view]"])

    # 7. Villas & Cottages
    build_content_rl(prs, "Accommodation — Private Villas & Lake Cottages",
        ["For the ultimate Finnish summer experience, our lakeside villas and "
         "cottages offer complete privacy on the water's edge — each with its "
         "own sauna, rowing boat, and direct lake access.",
         "From luxury villas with modern amenities to traditional log cottages "
         "deep in the forest, there's a perfect retreat for every taste.",
         "Many cottages include their own private island or lakefront, "
         "offering a truly exclusive Finnish escape."],
        ["[PHOTO: Lakeside villa exterior — summer]", "[PHOTO: Cottage with sauna & lake dock]"])

    # 8. Unique stays
    build_content_rl(prs, "Accommodation — Unique Experiences",
        ["The Finnish Lakeland offers unforgettable unique stays: floating saunas, "
         "treehouses overlooking the lake, and glass-walled aurora cabins that "
         "double as midnight sun viewpoints in summer.",
         "Järvisydän's Kota Hotel offers traditional Finnish kota-style accommodation "
         "— a modern take on the Lappish hut, right on the lakeshore.",
         "These one-of-a-kind stays create memories that last a lifetime."],
        ["[PHOTO: Kota Hotel or floating sauna]", "[PHOTO: Treehouse or unique accommodation]"])

    # 9. Product examples
    build_product(prs, "Build Your Lakeland Experience",
        "Top things to do — your clients choose how many days and what they want",
        [(2, "Weekend escape.\nLake cruise + sauna\n+ local dining\n+ hotel/cottage"),
         (3, "Classic Lakeland.\nSeal watching +\nmidnight fishing +\nsmoke sauna ritual"),
         (4, "Deep nature.\nAdd hiking,\nberry picking,\nfarm visit + cooking"),
         (5, "Multi-destination.\nCombine Järvisydän +\nTahko or add\nSahalahti retreat"),
         (6, "Full experience.\nTwo destinations +\nprivate island night\n+ all highlights"),
         (7, "Complete Lakeland.\nThree destinations,\nevery experience,\nmidnight sun magic")])

    # 10. Closing
    build_closing(prs, "All of this and much more\n— Lakeland awaits!",
                  "[CLOSING: Full-bleed midnight sun over lake]")

    # 11. Company
    build_company(prs)

    path = os.path.join(OUTPUT_DIR, "04-Lakeland-Summer.pptx")
    prs.save(path)
    print(f"  Saved: {path}")
    return path


# ========================================================================
# PRESENTATION 5: HELSINKI — Capital City (11 pages)
# ========================================================================
def build_05_helsinki():
    prs = new_prs()

    # 1. Cover
    build_cover(prs, "Helsinki", "The Nordic Urban Experience\nin the Finnish Capital",
                "[COVER: Helsinki waterfront panorama — Cathedral, harbour, sunset]")

    # 2. Stats
    build_stats(prs, "Welcome to Helsinki!", [
        ("1.5M", "metro area population"),
        ("300+", "islands in the archipelago"),
        ("2,000+", "restaurants in the city"),
        ("#1", "World Design Capital heritage"),
        ("Year-round", "something happening every season"),
    ], "[PHOTO: Helsinki aerial — harbour, cathedral, or design district]")

    # 3. Nordic Design Capital
    build_content_lr(prs, "Nordic Design Capital",
        ["Helsinki is a city that punches far above its weight. Compact and "
         "walkable, it combines world-class architecture, cutting-edge design, "
         "and a vibrant food scene into an unforgettable urban experience.",
         "From the neoclassical Helsinki Cathedral to the futuristic Oodi Library, "
         "from the historic Market Square to the trendy Design District — Helsinki "
         "is a city of contrasts and surprises.",
         "The archipelago of over 300 islands adds a unique coastal dimension, "
         "with island-hopping tours, seaside saunas, and waterfront dining "
         "just minutes from the city centre."],
        "[PHOTO: Helsinki Design District or Oodi Library]")

    # 4. Highlights
    build_destinations(prs, "Helsinki Highlights",
        "A compact city packed with world-class experiences",
        [("Architecture\n& Design", "Nordic Creativity",
          "Helsinki Cathedral\nOodi Library\nKamppi Chapel of Silence\nDesign District\nAalto & Saarinen landmarks"),
         ("Food &\nMarkets", "Nordic Gastronomy",
          "Old Market Hall (1889)\nHakaniemi Market\nRestaurant Day culture\nNew Nordic cuisine\nStreet food scene"),
         ("Islands &\nNature", "Urban Archipelago",
          "Suomenlinna Sea Fortress\n(UNESCO World Heritage)\nSeurasaari Open-Air Museum\nLonna Island\nPublic saunas")])

    # 5. Activities
    build_activities(prs, "Things to do in Helsinki!",
        "To offer the best possible experience, the city has endless options:",
        ["Suomenlinna Sea Fortress (UNESCO)", "Helsinki Cathedral & Senate Square",
         "Design District walking tour", "Old Market Hall food tour",
         "Archipelago island hopping", "Public sauna culture experience",
         "Oodi Library & Kamppi Chapel", "Temppeliaukio Rock Church",
         "Finnish National Museum", "Allas Sea Pool & sauna",
         "Craft beer & Nordic dining tour", "Street art & gallery walks",
         "Seasonal markets (Christmas!)", "Day trip to Tallinn or Porvoo"],
        ["[Helsinki Cathedral]", "[Market Hall]", "[Suomenlinna]",
         "[Design District]", "[Allas Sea Pool]", "[Archipelago]"])

    # 6. Hotels — City Centre
    build_content_rl(prs, "Accommodation — City Centre Hotels",
        ["Helsinki offers a wide range of centrally located hotels — from "
         "historic grande dames to sleek Scandinavian design hotels.",
         "Most city centre hotels are within walking distance of all major "
         "attractions, making Helsinki one of the easiest capitals to explore on foot.",
         "We partner with hotels across all categories to match your clients' "
         "style and budget."],
        ["[PHOTO: Helsinki hotel — exterior or lobby]", "[PHOTO: Hotel room with city view]"])

    # 7. Boutique & Design Hotels
    build_content_rl(prs, "Accommodation — Boutique & Design Hotels",
        ["Helsinki's design heritage extends to its hotels. Award-winning boutique "
         "properties showcase Finnish design, local art, and Nordic minimalism.",
         "From converted industrial spaces to waterfront properties with "
         "archipelago views, Helsinki's boutique hotels are destinations "
         "in themselves.",
         "These properties offer an authentic Helsinki experience that goes "
         "beyond just a place to sleep."],
        ["[PHOTO: Boutique hotel — design interior]", "[PHOTO: Waterfront or design hotel]"])

    # 8. Unique stays
    build_content_rl(prs, "Accommodation — Unique Helsinki Stays",
        ["Beyond traditional hotels, Helsinki offers unique accommodation "
         "experiences: island hotels on the archipelago, converted ships "
         "in the harbour, and floating saunas with overnight options.",
         "For groups, Helsinki has exceptional meeting and event venues "
         "in historic settings — from sea fortresses to art nouveau palaces.",
         "These unique options make Helsinki stand out as more than just "
         "a transit hub."],
        ["[PHOTO: Unique Helsinki stay — island or ship hotel]",
         "[PHOTO: Event venue or unique accommodation]"])

    # 9. Product examples
    build_product(prs, "Build Your Helsinki Experience",
        "Top things to do — your clients choose how many days and what they want",
        [(1, "City highlights.\nCathedral + Design\nDistrict + Market\nHall + dinner"),
         (2, "Classic Helsinki.\nAdd Suomenlinna\n+ sauna + archipelago\n+ evening dining"),
         (3, "Deep city.\nAdd museum day,\nfood tour, island\nhopping + culture"),
         (4, "Helsinki +\nDay trip to Porvoo\nor Tallinn.\nFull city experience"),
         (5, "Extended stay.\nAdd cooking class,\ngallery walks,\nlocal life immersion"),
         (6, "Helsinki +\nLakeland combo.\nCity + nature =\nperfect Finland intro"),
         (7, "Full week.\nHelsinki + day trips\n+ all highlights.\nThe complete Nordic\nurban experience")])

    # 10. Closing
    build_closing(prs, "All of this and much more\n— Helsinki awaits!",
                  "[CLOSING: Full-bleed Helsinki waterfront at sunset or winter lights]")

    # 11. Company
    build_company(prs)

    path = os.path.join(OUTPUT_DIR, "05-Helsinki-Capital.pptx")
    prs.save(path)
    print(f"  Saved: {path}")
    return path


# ========================================================================
# MAIN
# ========================================================================
if __name__ == "__main__":
    print("Building Finland DMC Presentations...")
    print("=" * 50)
    build_01_finland()
    build_02_dmc()
    build_03_lapland()
    build_04_lakeland()
    build_05_helsinki()
    print("=" * 50)
    print("All 5 presentations built!")
    print(f"Output directory: {OUTPUT_DIR}")
    print("\nNext steps:")
    print("  1. Open each .pptx in PowerPoint/Keynote")
    print("  2. Replace [PHOTO] placeholders with real DMC photos")
    print("  3. Adjust fonts if Playfair Display / Montserrat not installed")
    print("  4. Export as PDF for sharing")
