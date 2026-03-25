#!/usr/bin/env python3
"""
Levi Tour Presentation — Close Copy
Replicates the exact structure and feel of "Introduction to Levi as a Destination.pdf"
but branded for Finland DMC.

Key design rules from Levi Tour:
1. Warm cream background (#F5F0EB) — NOT white
2. Elegant serif headlines (Georgia/Playfair Display) — large, refined
3. Very generous whitespace — slides breathe
4. Photo occupies ~40-50% of slide area
5. Body text is light weight, dark gray, well-spaced
6. Minimal decoration — no boxes, no badges, no icons
7. Stats use oversized serif numbers
8. Cover + closing = full-bleed photo with overlaid text
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ========== DESIGN CONSTANTS (from Levi Tour) ==========
CREAM_BG = RGBColor(0xF5, 0xF0, 0xEB)
DARK_TEXT = RGBColor(0x2A, 0x2A, 0x2A)
MEDIUM_TEXT = RGBColor(0x4A, 0x4A, 0x4A)
LIGHT_TEXT = RGBColor(0x7A, 0x7A, 0x7A)
COVER_DARK = RGBColor(0x1A, 0x26, 0x34)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BRAND_BLUE = RGBColor(0x00, 0x35, 0xA0)
PHOTO_BG = RGBColor(0xDA, 0xD2, 0xC8)  # Muted warm tone for photo placeholders

# Fonts — Georgia is universally available and close to Playfair Display feel
SERIF = 'Georgia'
SANS = 'Gill Sans'  # Clean sans-serif available on macOS

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

LOGO_PATH = os.path.expanduser("~/Downloads/Media/Fin+Land-dmc.png")
OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, left, top, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def photo_box(slide, left, top, w, h, label=""):
    """Warm-toned photo placeholder with centered label."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = PHOTO_BG
    s.line.fill.background()
    tf = s.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Vertical center
    s.text_frame.paragraphs[0].space_before = Pt(0)
    run = tf.paragraphs[0].add_run()
    run.text = label
    run.font.name = SANS
    run.font.size = Pt(11)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x9A, 0x90, 0x85)
    return s

def txt(slide, text, left, top, w, h, font=SERIF, size=18, bold=False, italic=False,
        color=DARK_TEXT, align=PP_ALIGN.LEFT, spacing=1.15):
    """Add a single text box."""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    p.line_spacing = Pt(int(size * spacing))
    return tb

def para_block(slide, paragraphs, left, top, w, h, size=14, color=MEDIUM_TEXT):
    """Multiple paragraphs in one text box."""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, text in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = SANS
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = Pt(int(size * 1.6))
        p.space_after = Pt(14)
    return tb

def logo_bottom_right(slide):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, SLIDE_W - Inches(3), SLIDE_H - Inches(0.7), height=Inches(0.45))


# ========================================================================
# SLIDE BUILDERS — exact Levi Tour layout
# ========================================================================

def slide_cover(prs, destination, tagline, photo_desc):
    """Slide 1: Full-bleed photo cover with destination name in cream band at bottom."""
    s = blank(prs)
    # Full photo background
    photo_box(s, 0, 0, SLIDE_W, Inches(5),
              f"[FULL-BLEED COVER PHOTO]\n{photo_desc}")
    # Cream band at bottom (like Levi Tour)
    rect(s, 0, Inches(4.5), SLIDE_W, Inches(3), CREAM_BG)
    # Destination name — large elegant serif
    txt(s, destination, Inches(1), Inches(4.8), Inches(11), Inches(1.2),
        font=SERIF, size=54, bold=False, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    # Tagline
    txt(s, tagline, Inches(1.5), Inches(6.0), Inches(10), Inches(0.8),
        font=SERIF, size=20, italic=True, color=MEDIUM_TEXT, align=PP_ALIGN.CENTER)
    return s

def slide_stats(prs, title, stats, photo_desc):
    """Slide 2: Welcome stats — photo left, big numbers right.
    stats = [(number, label), ...]"""
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, CREAM_BG)
    # Photo left ~35%
    photo_box(s, Inches(0.5), Inches(0.5), Inches(4.5), Inches(6.5), photo_desc)
    # Title
    txt(s, title, Inches(5.5), Inches(0.8), Inches(7), Inches(0.8),
        font=SERIF, size=36, color=DARK_TEXT)
    # Stats in grid — oversized numbers like Levi Tour
    positions = [
        (Inches(5.5), Inches(2.0)),
        (Inches(9.0), Inches(2.0)),
        (Inches(5.5), Inches(3.8)),
        (Inches(9.0), Inches(3.8)),
    ]
    for i, (num, label) in enumerate(stats[:4]):
        if i < len(positions):
            x, y = positions[i]
            # Big number
            txt(s, num, x, y, Inches(3.2), Inches(0.8),
                font=SERIF, size=44, bold=True, color=DARK_TEXT)
            # Small label
            txt(s, label, x, y + Inches(0.7), Inches(3.2), Inches(0.5),
                font=SANS, size=13, color=LIGHT_TEXT)
    # Bottom center stat if 5th exists
    if len(stats) > 4:
        num, label = stats[4]
        txt(s, num, Inches(7), Inches(5.5), Inches(4), Inches(0.8),
            font=SERIF, size=44, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        txt(s, label, Inches(7), Inches(6.2), Inches(4), Inches(0.4),
            font=SANS, size=13, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)
    logo_bottom_right(s)
    return s

def slide_text_photo(prs, title, paragraphs, photo_desc):
    """Slide 3: Text left, big photo right (like Northern Lights slide)."""
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, CREAM_BG)
    # Title
    txt(s, title, Inches(0.8), Inches(0.8), Inches(6), Inches(1),
        font=SERIF, size=34, color=DARK_TEXT)
    # Body text
    para_block(s, paragraphs, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5))
    # Photo right
    photo_box(s, Inches(7), Inches(0.3), Inches(6), Inches(6.9), photo_desc)
    logo_bottom_right(s)
    return s

def slide_activities(prs, title, subtitle, activities, photo_descs):
    """Slide 4: Activities — text left, photo mosaic right (2x3 grid)."""
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, CREAM_BG)
    # Title
    txt(s, title, Inches(0.8), Inches(0.5), Inches(5.5), Inches(1.2),
        font=SERIF, size=30, color=DARK_TEXT)
    # Subtitle
    txt(s, subtitle, Inches(0.8), Inches(1.6), Inches(5), Inches(0.5),
        font=SANS, size=12, italic=True, color=MEDIUM_TEXT)
    # Activity list — 2 columns
    mid = len(activities) // 2
    col1 = activities[:mid]
    col2 = activities[mid:]
    # Column 1
    tb1 = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(2.8), Inches(4.5)) if False else None
    tb1_box = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(2.8), Inches(4.5))
    tf1 = tb1_box.text_frame
    tf1.word_wrap = True
    for i, act in enumerate(col1):
        p = tf1.paragraphs[0] if i == 0 else tf1.add_paragraph()
        p.text = f"\u2022  {act}"
        p.font.name = SANS
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)
    # Column 2
    tb2_box = s.shapes.add_textbox(Inches(3.6), Inches(2.4), Inches(3), Inches(4.5))
    tf2 = tb2_box.text_frame
    tf2.word_wrap = True
    for i, act in enumerate(col2):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"\u2022  {act}"
        p.font.name = SANS
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)
    # "...and many more!"
    txt(s, "...and many more!", Inches(0.8), Inches(6.5), Inches(3), Inches(0.5),
        font=SANS, size=12, italic=True, color=LIGHT_TEXT)
    # Photo mosaic — 2 cols x 3 rows on right side
    pw = Inches(2.85)
    ph = Inches(2.3)
    gap = Inches(0.1)
    for i in range(6):
        col = i % 2
        row = i // 2
        x = Inches(7.2) + col * (pw + gap)
        y = Inches(0.3) + row * (ph + gap)
        label = photo_descs[i] if i < len(photo_descs) else f"[Activity photo {i+1}]"
        photo_box(s, x, y, pw, ph, label)
    logo_bottom_right(s)
    return s

def slide_accommodation(prs, title, paragraphs, photo_descs):
    """Slides 5-7: 2 stacked photos left, text right (exact Levi Tour layout)."""
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, CREAM_BG)
    # 2 stacked photos left
    photo_box(s, Inches(0.5), Inches(0.3), Inches(5), Inches(3.4),
              photo_descs[0] if photo_descs else "[Photo 1]")
    photo_box(s, Inches(0.5), Inches(3.8), Inches(5), Inches(3.4),
              photo_descs[1] if len(photo_descs) > 1 else "[Photo 2]")
    # Title right
    txt(s, title, Inches(6), Inches(1), Inches(6.8), Inches(1),
        font=SERIF, size=32, color=DARK_TEXT)
    # Body text
    para_block(s, paragraphs, Inches(6), Inches(2.5), Inches(6.5), Inches(4))
    logo_bottom_right(s)
    return s

def slide_closing(prs, text, photo_desc):
    """Slide 8: Full-bleed atmospheric photo + centered text."""
    s = blank(prs)
    photo_box(s, 0, 0, SLIDE_W, SLIDE_H, photo_desc)
    # Dark overlay for text readability
    overlay = rect(s, 0, 0, SLIDE_W, SLIDE_H, COVER_DARK)
    # Can't set transparency in python-pptx, so use dark bg
    # Text centered
    txt(s, text, Inches(1.5), Inches(2.5), Inches(10), Inches(2.5),
        font=SERIF, size=36, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
    logo_bottom_right(s)
    return s

def slide_company(prs, company_name, tagline, description_paras, phone, email, web):
    """Slide 9: Company page — logo centered, description, contact."""
    s = blank(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, CREAM_BG)
    # Logo
    if os.path.exists(LOGO_PATH):
        s.shapes.add_picture(LOGO_PATH, Inches(5), Inches(0.5), height=Inches(1))
    # Tagline
    txt(s, tagline, Inches(2), Inches(2), Inches(9), Inches(0.6),
        font=SERIF, size=18, italic=True, color=MEDIUM_TEXT, align=PP_ALIGN.CENTER)
    # Description
    tb = s.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, para in enumerate(description_paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = para
        p.font.name = SANS
        p.font.size = Pt(12)
        p.font.color.rgb = MEDIUM_TEXT
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = Pt(20)
        p.space_after = Pt(12)
    # Contact
    txt(s, web, Inches(3), Inches(6.0), Inches(7), Inches(0.3),
        font=SANS, size=13, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    txt(s, f"{phone}   |   {email}", Inches(3), Inches(6.4), Inches(7), Inches(0.3),
        font=SANS, size=12, color=MEDIUM_TEXT, align=PP_ALIGN.CENTER)
    return s


# ========================================================================
# BUILD: Levi Tour Copy (Introduction to Levi as a Destination)
# ========================================================================
def build_levi_copy():
    prs = new_prs()

    # 1. Cover
    slide_cover(prs, "Levi",
        "The Home of The Northern Lights\nin Finnish Lapland",
        "[Northern Lights over snow-covered trees — iconic Levi landscape]")

    # 2. Welcome + Stats
    slide_stats(prs, "Welcome to Levi!", [
        ("#1", "largest ski resort\nin Lapland"),
        ("111", "yearly Northern\nLights sightings"),
        ("60", "restaurants\nand bars"),
        ("25,000", "beds for\naccommodation"),
        ("750,000", "visitors\na year"),
    ], "[Levi fell at sunset —\naerial view with ski slopes]")

    # 3. Northern Lights
    slide_text_photo(prs, "The Northern Lights", [
        "Levi is one of the best places in the world to experience the Northern Lights.",
        "Located far above the Arctic Circle and directly under the auroral oval, Levi offers excellent chances to see the Aurora Borealis on clear nights.",
        "Thanks to long polar nights, clean Arctic air, and very low light pollution, the skies around Levi are ideal for aurora watching.",
    ], "[Aurora Borealis over\nsnow-covered landscape near Levi]")

    # 4. Activities
    slide_activities(prs,
        "Things to do in the\nunforgettable Lappish settings!",
        "To offer the best possible experience for all the family, in Levi we have plenty of things to do and explore:",
        ["Snowmobile safari", "Northern Lights sighting", "Snow Castle visit",
         "Snowshoeing", "Icekarting", "Husky & Reindeer Safaris",
         "Cross country skiing", "Alpine skiing", "Spa & Water World",
         "Day spas & Wellness", "Helicopter & Snowcat tours", "Meeting Santa Claus"],
        ["[Husky safari]", "[Snowmobile safari]",
         "[Reindeer ride — sunset]", "[Ice karting]",
         "[Cross country skiing\nin snowy forest]", "[Snowcat tour\nor helicopter view]"])

    # 5. Hotels
    slide_accommodation(prs, "Accommodation –\nHotels", [
        "The range of hotels in Levi offer a perfect balance of Arctic atmosphere and modern comfort. From cosy alpine hotels to stylish design stays and family-friendly spa hotel, there's something for every type of traveller.",
        "Many hotels are located close to the slopes, activities, and village services ensuring easy access.",
    ], ["[Levi hotel exterior —\nwinter, snow-covered]",
        "[Levi hotel — modern design\nor lobby/restaurant view]"])

    # 6. Private Chalets
    slide_accommodation(prs, "Accommodation –\nPrivate Chalets", [
        "Levi offers private chalet accommodation from bigger private villas to smaller apartments in the Alpine chalets.",
        "Whether you are looking for complete privacy or easy access to the centre of the village, there is something for everyone.",
        "This accommodation type provides more flexibility to the stay and can be tailored to the customer's needs.",
    ], ["[Private chalet exterior —\nAlpine style, winter]",
        "[Chalet or villa —\nsnowy forest setting]"])

    # 7. Northern Lights Igloos
    slide_accommodation(prs, "Accommodation –\nNorthern Lights Igloos", [
        "The Aurora Cabins and Northern Light Igloos are a magical way to experience the wintry star-filled sky and the Northern Lights, while lying comfortably in the warm bed.",
        "Levi has multiple resorts that are located further away from the light pollution offering a peaceful stay and optimal setting for Northern lights travelling.",
    ], ["[Glass igloo village —\naerial view, snowy forest]",
        "[Inside glass igloo —\naurora visible through ceiling]"])

    # 8. Closing
    slide_closing(prs,
        "All of this and much more\n— Lapland awaits!",
        "[Full-bleed Northern Lights —\nvivid green aurora over Levi landscape]")

    # 9. Company
    slide_company(prs,
        "Finland DMC",
        "We know hospitality.",
        ["Finland DMC is a specialist destination management company covering all of Finland. For over a decade, we have been crafting beautifully detailed holidays that combine authentic Arctic nature with comfort, quality, and local expertise.",
         "We operate at the very heart of our guests' journeys, offering seamless, end-to-end travel experiences. From carefully selected accommodation and dining to inspiring activities and excursions, we bring every element together into one smooth and memorable whole.",
         "Our services are designed to remove complexity and allow our guests to fully enjoy the magic of Finland. We serve a wide range of travellers — from couples and families to larger groups, corporate teams, and event organisers.",
         "With deep local knowledge, long-term partnerships, and a genuine passion for Finland, we turn trips into meaningful experiences that stay with our guests long after they return home."],
        "+358 43 2010 687",
        "info@finland-dmc.com",
        "www.finland-dmc.com")

    path = os.path.join(OUTPUT_DIR, "00-Levi-Tour-Copy.pptx")
    prs.save(path)
    print(f"Saved: {path}")
    return path


if __name__ == "__main__":
    print("Building Levi Tour Copy...")
    build_levi_copy()
    print("Done! Open in PowerPoint/Keynote to review.")
